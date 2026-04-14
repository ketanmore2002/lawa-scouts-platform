"""E2B sandbox for executing Python code — charts, PDF/Excel exports."""

import json
import base64
import logging
import re
from app.config import get_settings

def _slugify(text: str, max_len: int = 50) -> str:
    """Create a filename-safe slug from text."""
    slug = re.sub(r'[^\w\s-]', '', text.lower().strip())
    slug = re.sub(r'[\s_]+', '-', slug)
    return slug[:max_len].rstrip('-') or 'report'

logger = logging.getLogger(__name__)

# Graceful import
_E2B_AVAILABLE = False
try:
    from e2b_code_interpreter import Sandbox
    _E2B_AVAILABLE = True
except ImportError:
    logger.info("e2b-code-interpreter not installed — sandbox disabled")


def is_available() -> bool:
    return _E2B_AVAILABLE and bool(get_settings().e2b_api_key)


# ────────────────────────────────────────────────────────
# Data analysis — should we even chart this?
# ────────────────────────────────────────────────────────

_NUMERIC_HINTS = re.compile(
    r"price|cost|rating|score|salary|revenue|count|amount|percent|%|"
    r"budget|funding|market.?cap|valuation|growth|index|rank|weight|"
    r"temperature|population|gdp|downloads|users|views|subscribers",
    re.IGNORECASE,
)


def _has_chartable_data(columns: list[dict], rows: list[dict]) -> bool:
    """Check if the structured data has numeric/categorical fields worth charting."""
    if len(rows) < 3:
        return False

    # Check column keys/labels for numeric hints
    for col in columns:
        key = col.get("key", "")
        label = col.get("label", "")
        if _NUMERIC_HINTS.search(key) or _NUMERIC_HINTS.search(label):
            return True

    # Check actual row values — are there numbers hiding in text fields?
    numeric_cols = 0
    for col in columns:
        if col.get("type") in ("link", "tags"):
            continue
        key = col.get("key", "")
        num_count = 0
        for row in rows[:10]:
            val = str(row.get(key, ""))
            cleaned = re.sub(r"[,$€£₹¥%\s]", "", val)
            try:
                float(cleaned)
                num_count += 1
            except (ValueError, TypeError):
                pass
        if num_count >= 3:
            numeric_cols += 1

    # Check badge columns: need 2+ distinct values AND not nearly-unique
    # (if >60% of values are unique, a distribution chart would be messy)
    useful_badge_cols = 0
    for col in columns:
        if col.get("type") == "badge":
            key = col.get("key", "")
            sample = [str(row.get(key, "")) for row in rows[:20] if row.get(key)]
            distinct = len(set(sample))
            total = len(sample)
            if distinct >= 2 and total > 0 and (distinct / total) <= 0.5:
                useful_badge_cols += 1

    if numeric_cols >= 1:
        return True
    if useful_badge_cols >= 1 and len(rows) >= 5:
        return True

    return False


# ────────────────────────────────────────────────────────
# Chart generation
# ────────────────────────────────────────────────────────

_CHART_PROMPT = """Write Python code to create 1-2 clean, BEAUTIFUL data visualisation charts for this dataset about "{topic}".

Data sample (JSON — full data is in the DATA variable):
{sample}

DESIGN RULES — follow these EXACTLY:
1. Use matplotlib with this exact setup at the top:
   import matplotlib
   matplotlib.use('Agg')
   import matplotlib.pyplot as plt
   import matplotlib.ticker as ticker
   from matplotlib import rcParams
   from datetime import datetime

2. Apply this EXACT theme before any plotting:
   rcParams['font.family'] = 'sans-serif'
   rcParams['font.sans-serif'] = ['DejaVu Sans']
   rcParams['font.size'] = 11
   rcParams['axes.spines.top'] = False
   rcParams['axes.spines.right'] = False
   rcParams['axes.linewidth'] = 0.6
   rcParams['axes.edgecolor'] = '#CBD5E1'
   rcParams['axes.labelcolor'] = '#475569'
   rcParams['xtick.color'] = '#64748B'
   rcParams['ytick.color'] = '#64748B'
   rcParams['figure.facecolor'] = '#FFFFFF'
   rcParams['axes.facecolor'] = '#FFFFFF'
   rcParams['grid.color'] = '#F1F5F9'
   rcParams['grid.linewidth'] = 0.5

3. Color palette — use ONLY these colors:
   COLORS = ['#0E7490', '#0891B2', '#22D3EE', '#67E8F9', '#A5F3FC',
             '#155E75', '#164E63', '#083344', '#06B6D4', '#00BCD4']

4. Chart style rules:
   - fig, ax = plt.subplots(figsize=(9, 5)) for each chart
   - Horizontal bar charts are preferred over vertical when labels are long
   - Truncate long labels to 30 chars with "..."
   - ALWAYS add a clear, descriptive caption below the chart using:
     fig.text(0.5, -0.02, 'Your caption here', ha='center', fontsize=10, color='#64748B', style='italic')
     The caption must describe what the CHART shows in plain English. Write it yourself based on what the chart actually visualizes.
     DO NOT copy or include the user's raw query in the caption. Extract the core subject and write a clean caption.
     For example, if the topic is "Track iran and uae war news along few charts and graphs", a good caption is "Top categories in Iran-UAE war news coverage" — NOT "Top news themes in Track iran and uae war news along few charts and graphs".
     NEVER use the words "dataset", "data points", "data", "sample", or "Track" in captions.
     NEVER paste the raw topic string into the caption. Rephrase it.
   - Subtle grid on the value axis ONLY: ax.grid(axis='x', alpha=0.3) or ax.grid(axis='y', alpha=0.3)
   - Add value labels on bars: ax.bar_label(bars, fmt='%.0f', padding=3, fontsize=9, color='#475569')
   - Rounded bar edges where possible: use edgecolor='white', linewidth=0.5
   - plt.tight_layout(pad=2.0)

5. PIE / DONUT CHART RULES (CRITICAL — follow exactly):
   - MAXIMUM 8 slices. If more than 8 categories, keep the top 7 by count and group ALL others into a single "Other" slice.
   - Use ax.pie with autopct ONLY when there are <= 8 slices: wedges, texts, autotexts = ax.pie(sizes, labels=labels, autopct='%1.0f%%', colors=COLORS, startangle=90, wedgeprops=dict(width=0.6, edgecolor='white'), pctdistance=0.75, textprops=dict(fontsize=10))
   - Set label font size small: for t in texts: t.set_fontsize(9)
   - For autotext: for t in autotexts: t.set_fontsize(8); t.set_color('#475569')
   - If there are too many categories for a donut to look good, use a HORIZONTAL BAR chart instead showing top 10.
   - NEVER create a pie/donut with more than 8 slices. It will look terrible.

6. TIME SERIES / LINE CHART RULES (CRITICAL — follow exactly):
   - Parse dates properly. Try multiple formats: '%Y-%m-%d', '%b %d, %Y', '%d/%m/%Y', '%Y-%m-%dT%H:%M:%S', etc.
   - Use matplotlib.dates for formatting: import matplotlib.dates as mdates
   - Choose date format based on span:
     * If data spans more than 365 days: use '%b %Y' (e.g., "Jan 2025")
     * If data spans 30-365 days: use '%b %d' (e.g., "Jan 15")
     * If data spans less than 30 days: use '%b %d' (e.g., "Mar 28")
     Calculate span: span_days = (max(dates) - min(dates)).days
     Then: fmt = '%b %Y' if span_days > 365 else '%b %d'
     ax.xaxis.set_major_formatter(mdates.DateFormatter(fmt))
   - Auto-locate ticks sensibly: ax.xaxis.set_major_locator(mdates.AutoDateLocator(minticks=4, maxticks=10))
   - Rotate date labels: plt.setp(ax.xaxis.get_majorticklabels(), rotation=30, ha='right')
   - Use area fill: ax.fill_between(dates, values, alpha=0.1, color=COLORS[0])
   - Use ax.yaxis.set_major_locator(ticker.MaxNLocator(integer=True)) for integer y-axes
   - Add proper axis labels: ax.set_xlabel('Date', fontsize=10), ax.set_ylabel('Count', fontsize=10)

7. Save EACH chart: plt.savefig('/tmp/chart_N.png', dpi=150, bbox_inches='tight', facecolor='white')
   then plt.close()

8. Only create charts that MAKE SENSE for this data:
   - Prices/costs/amounts → horizontal bar comparison of top 10
   - Ratings/scores → sorted bar chart
   - Categories with <= 8 unique values → donut chart; otherwise → horizontal bar of top 10
   - Time series → clean line chart with proper date axis
   - If data is purely textual (just names, descriptions, links) → print("NO_CHARTS") and create NO charts

9. Handle missing/None values: skip rows with missing values for the charted field.

10. Print a 2-3 line analysis summary to stdout.

Return ONLY Python code. No markdown fences."""


def run_analysis(
    data: list[dict],
    scout_topic: str,
    columns: list[dict] | None = None,
    force_charts: bool = False,
) -> dict:
    """Generate and execute analysis code in E2B sandbox.

    Args:
        data: Row dicts from the structured report.
        scout_topic: Topic for context.
        columns: Column definitions (used to check chartability).
        force_charts: If True, skip chartability check (user explicitly asked for charts).

    Returns:
        dict with keys: analysis_text, charts (list of base64 PNGs), error
    """
    if not is_available():
        return {"analysis_text": "", "charts": [], "error": "E2B not available"}

    # Gate: only chart when data has chartable fields (unless user asked for charts)
    if not force_charts and columns and not _has_chartable_data(columns, data):
        logger.info("Data not chartable, skipping sandbox analysis")
        return {"analysis_text": "", "charts": [], "error": None}

    try:
        from openai import OpenAI
        client = OpenAI(api_key=get_settings().openai_api_key)

        data_sample = data[:15]
        prompt = _CHART_PROMPT.format(
            topic=scout_topic,
            sample=json.dumps(data_sample, indent=2),
        )

        code_resp = client.responses.create(model="gpt-5.4-mini-2026-03-17", input=prompt)

        code = code_resp.output_text.strip()
        if code.startswith("```"):
            code = code.split("\n", 1)[1] if "\n" in code else code[3:]
        if code.endswith("```"):
            code = code[:-3]

        full_code = f"DATA = {json.dumps(data)}\n\n{code}"

        logger.info(f"Sandbox: executing {len(full_code)} chars of analysis code")
        return _execute_in_sandbox(full_code)

    except Exception as e:
        logger.error(f"Sandbox analysis failed: {e}")
        return {"analysis_text": "", "charts": [], "error": str(e)}


# ────────────────────────────────────────────────────────
# Multi-format export templates
# ────────────────────────────────────────────────────────

_PDF_CODE = '''
import json, textwrap

DATA = {data_json}
TITLE = {title_json}
SUMMARY = {summary_json}
COLUMNS = {columns_json}
TOPIC = {topic_json}
ANALYSIS = {analysis_json}

from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib import colors
from reportlab.lib.units import mm
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, HRFlowable
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

doc = SimpleDocTemplate("/tmp/report.pdf", pagesize=landscape(A4),
                        leftMargin=15*mm, rightMargin=15*mm,
                        topMargin=15*mm, bottomMargin=15*mm)
styles = getSampleStyleSheet()

title_style = ParagraphStyle('CustomTitle', parent=styles['Title'],
                              fontSize=18, textColor=colors.HexColor('#0E7490'),
                              spaceAfter=6)
subtitle_style = ParagraphStyle('Subtitle', parent=styles['Normal'],
                                 fontSize=10, textColor=colors.HexColor('#64748B'),
                                 spaceAfter=16)
body_style = ParagraphStyle('Body', parent=styles['Normal'],
                             fontSize=9, leading=13,
                             textColor=colors.HexColor('#334155'))
analysis_heading_style = ParagraphStyle('AnalysisHeading', parent=styles['Normal'],
                             fontSize=12, leading=16, fontName='Helvetica-Bold',
                             textColor=colors.HexColor('#0E7490'),
                             spaceBefore=8, spaceAfter=6)
analysis_style = ParagraphStyle('Analysis', parent=styles['Normal'],
                             fontSize=9, leading=14,
                             textColor=colors.HexColor('#334155'),
                             spaceAfter=8)
cell_style = ParagraphStyle('Cell', parent=styles['Normal'],
                             fontSize=8, leading=10,
                             textColor=colors.HexColor('#334155'))

elements = []
elements.append(Paragraph(TITLE, title_style))
elements.append(Paragraph("LAWA Scouts Report &mdash; " + TOPIC, subtitle_style))
elements.append(Paragraph(SUMMARY, body_style))
elements.append(Spacer(1, 8))

# Analysis section
if ANALYSIS and len(ANALYSIS.strip()) > 20:
    elements.append(HRFlowable(width="100%", thickness=0.5, color=colors.HexColor('#E2E8F0'), spaceAfter=10, spaceBefore=4))
    elements.append(Paragraph("Analysis &amp; Key Findings", analysis_heading_style))
    for para in ANALYSIS.strip().split("\\n\\n"):
        para = para.strip()
        if para:
            # Clean up markdown bold markers for PDF
            para = para.replace("**", "")
            elements.append(Paragraph(para, analysis_style))
    elements.append(HRFlowable(width="100%", thickness=0.5, color=colors.HexColor('#E2E8F0'), spaceAfter=10, spaceBefore=4))

elements.append(Spacer(1, 8))
elements.append(Paragraph("<b>Detailed Data</b>", analysis_heading_style))
elements.append(Spacer(1, 4))

col_keys = [c["key"] for c in COLUMNS if c.get("type") != "link"]
col_labels = [c["label"] for c in COLUMNS if c.get("type") != "link"]

header = [Paragraph("<b>" + l + "</b>", cell_style) for l in col_labels]
table_data = [header]
for row in DATA:
    cells = []
    for k in col_keys:
        val = row.get(k, "")
        if isinstance(val, list):
            val = ", ".join(str(v) for v in val)
        cells.append(Paragraph(str(val or "\\u2014")[:120], cell_style))
    table_data.append(cells)

ncols = len(col_keys)
avail = 267 * mm  # landscape A4 minus margins
col_widths = [avail / ncols] * ncols

t = Table(table_data, colWidths=col_widths, repeatRows=1)
t.setStyle(TableStyle([
    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#F0FDFA')),
    ('TEXTCOLOR', (0, 0), (-1, 0), colors.HexColor('#0E7490')),
    ('FONTSIZE', (0, 0), (-1, -1), 8),
    ('GRID', (0, 0), (-1, -1), 0.4, colors.HexColor('#E2E8F0')),
    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#F8FAFC')]),
    ('VALIGN', (0, 0), (-1, -1), 'TOP'),
    ('TOPPADDING', (0, 0), (-1, -1), 4),
    ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
    ('LEFTPADDING', (0, 0), (-1, -1), 6),
    ('RIGHTPADDING', (0, 0), (-1, -1), 6),
]))
elements.append(t)

doc.build(elements)
print("PDF_GENERATED:/tmp/report.pdf")
'''

_EXCEL_CODE = '''
import json

DATA = {data_json}
TITLE = {title_json}
COLUMNS = {columns_json}

try:
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
except ImportError:
    import subprocess
    subprocess.check_call(["pip", "install", "openpyxl"])
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

wb = openpyxl.Workbook()
ws = wb.active
ws.title = "Report"

header_font = Font(name='Calibri', bold=True, size=10, color='0E7490')
header_fill = PatternFill(start_color='F0FDFA', end_color='F0FDFA', fill_type='solid')
cell_font = Font(name='Calibri', size=10, color='334155')
thin_border = Border(
    bottom=Side(style='thin', color='E2E8F0'),
    right=Side(style='thin', color='E2E8F0'),
)

ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=len(COLUMNS))
ws.cell(row=1, column=1, value=TITLE).font = Font(name='Calibri', bold=True, size=14, color='0E7490')
ws.row_dimensions[1].height = 30

col_keys = [c["key"] for c in COLUMNS]
for ci, col in enumerate(COLUMNS, 1):
    cell = ws.cell(row=3, column=ci, value=col["label"])
    cell.font = header_font
    cell.fill = header_fill
    cell.alignment = Alignment(horizontal='left', vertical='center')
    cell.border = thin_border

for ri, row in enumerate(DATA, 4):
    for ci, key in enumerate(col_keys, 1):
        val = row.get(key, "")
        if isinstance(val, list):
            val = ", ".join(str(v) for v in val)
        cell = ws.cell(row=ri, column=ci, value=str(val or ""))
        cell.font = cell_font
        cell.border = thin_border
        cell.alignment = Alignment(vertical='top', wrap_text=True)

for ci in range(1, len(COLUMNS) + 1):
    max_len = max(
        len(str(ws.cell(row=r, column=ci).value or ""))
        for r in range(3, min(len(DATA) + 4, 50))
    )
    ws.column_dimensions[openpyxl.utils.get_column_letter(ci)].width = min(40, max(12, max_len + 2))

wb.save("/tmp/report.xlsx")
print("EXCEL_GENERATED:/tmp/report.xlsx")
'''

_HTML_CODE = '''
import json, html as html_mod

DATA = {data_json}
TITLE = {title_json}
SUMMARY = {summary_json}
COLUMNS = {columns_json}
ANALYSIS = {analysis_json}

col_keys = [c["key"] for c in COLUMNS]
col_labels = [c["label"] for c in COLUMNS]

rows_html = ""
for row in DATA:
    cells = ""
    for c in COLUMNS:
        k = c["key"]
        val = row.get(k, "")
        if isinstance(val, list):
            val = ", ".join(str(v) for v in val)
        val = str(val or "\\u2014")
        if c.get("type") == "link" and val.startswith("http"):
            label = row.get("link_label", "Visit")
            cells += f\'<td><a href="{{val}}" target="_blank" style="color:#0E7490;text-decoration:underline;">{{label}}</a></td>\'
        elif c.get("type") == "badge":
            cells += f\'<td><span style="display:inline-block;padding:2px 10px;border-radius:999px;background:#ECFEFF;color:#0E7490;font-size:12px;font-weight:600;">{{val}}</span></td>\'
        else:
            cells += f"<td>{{val}}</td>"
    rows_html += f"<tr>{{cells}}</tr>\\n"

headers = "".join(f"<th>{{l}}</th>" for l in col_labels)

# Build analysis HTML
analysis_html = ""
if ANALYSIS and len(ANALYSIS.strip()) > 20:
    paras = ANALYSIS.strip().split("\\n\\n")
    analysis_parts = []
    for p in paras:
        p = p.strip()
        if p:
            # Convert **bold** to <strong>
            import re
            p = re.sub(r\'\\*\\*(.+?)\\*\\*\', r\'<strong>\\1</strong>\', html_mod.escape(p))
            p = p.replace("\\n", "<br>")
            analysis_parts.append(f"<p>{{p}}</p>")
    if analysis_parts:
        analysis_html = \'<div class="analysis"><h2>Analysis &amp; Key Findings</h2>\' + "".join(analysis_parts) + "</div>"

html_out = """<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>""" + html_mod.escape(TITLE) + """</title>
<style>
  *{margin:0;padding:0;box-sizing:border-box}
  body{font-family:'Inter','Segoe UI',system-ui,sans-serif;background:#F8FAFC;color:#334155;padding:40px 24px}
  .container{max-width:1100px;margin:0 auto}
  h1{font-size:24px;font-weight:700;color:#0E7490;margin-bottom:8px}
  .subtitle{font-size:13px;color:#64748B;margin-bottom:24px}
  .summary{font-size:14px;line-height:1.7;color:#475569;margin-bottom:24px;padding:16px 20px;background:#fff;border-radius:12px;border:1px solid #E2E8F0}
  .analysis{margin-bottom:32px;padding:20px 24px;background:#fff;border-radius:12px;border:1px solid #E2E8F0}
  .analysis h2{font-size:16px;font-weight:700;color:#0E7490;margin-bottom:14px}
  .analysis p{font-size:14px;line-height:1.75;color:#475569;margin-bottom:12px}
  .analysis p:last-child{margin-bottom:0}
  h3{font-size:14px;font-weight:600;color:#0E7490;margin:24px 0 8px}
  table{width:100%;border-collapse:collapse;background:#fff;border-radius:12px;overflow:hidden;border:1px solid #E2E8F0}
  th{background:#F0FDFA;color:#0E7490;font-size:12px;font-weight:600;text-transform:uppercase;letter-spacing:.5px;padding:12px 16px;text-align:left;border-bottom:2px solid #E2E8F0}
  td{padding:10px 16px;font-size:13px;border-bottom:1px solid #F1F5F9;vertical-align:top}
  tr:hover td{background:#F8FAFC}
  .footer{margin-top:32px;text-align:center;font-size:11px;color:#94A3B8}
</style>
</head>
<body>
<div class="container">
  <h1>""" + html_mod.escape(TITLE) + """</h1>
  <p class="subtitle">LAWA Scouts Report</p>
  <div class="summary">""" + html_mod.escape(SUMMARY) + """</div>
  """ + analysis_html + """
  <h3>Detailed Data</h3>
  <table>
    <thead><tr>""" + headers + """</tr></thead>
    <tbody>""" + rows_html + """</tbody>
  </table>
  <p class="footer">Generated by LAWA Scouts</p>
</div>
</body>
</html>"""

with open("/tmp/report.html", "w", encoding="utf-8") as f:
    f.write(html_out)
print("HTML_GENERATED:/tmp/report.html")
'''

_CSV_CODE = '''
import json, csv

DATA = {data_json}
COLUMNS = {columns_json}

col_keys = [c["key"] for c in COLUMNS]
col_labels = [c["label"] for c in COLUMNS]

with open("/tmp/report.csv", "w", newline="", encoding="utf-8") as f:
    writer = csv.writer(f)
    writer.writerow(col_labels)
    for row in DATA:
        cells = []
        for k in col_keys:
            val = row.get(k, "")
            if isinstance(val, list):
                val = ", ".join(str(v) for v in val)
            cells.append(str(val or ""))
        writer.writerow(cells)

print("CSV_GENERATED:/tmp/report.csv")
'''

_TXT_CODE = '''
import json, textwrap

DATA = {data_json}
TITLE = {title_json}
SUMMARY = {summary_json}
COLUMNS = {columns_json}
ANALYSIS = {analysis_json}

col_keys = [c["key"] for c in COLUMNS if c.get("type") != "link"]
col_labels = [c["label"] for c in COLUMNS if c.get("type") != "link"]

lines = []
lines.append("=" * 80)
lines.append(TITLE)
lines.append("=" * 80)
lines.append("")
lines.append(textwrap.fill(SUMMARY, width=80))
lines.append("")

# Analysis section
if ANALYSIS and len(ANALYSIS.strip()) > 20:
    lines.append("-" * 80)
    lines.append("ANALYSIS & KEY FINDINGS")
    lines.append("-" * 80)
    lines.append("")
    for para in ANALYSIS.strip().split("\\n\\n"):
        para = para.strip().replace("**", "")
        if para:
            lines.append(textwrap.fill(para, width=80))
            lines.append("")

lines.append("-" * 80)
lines.append("DETAILED DATA")
lines.append("-" * 80)

for ri, row in enumerate(DATA, 1):
    lines.append(f"\\n[{{ri}}]")
    for k, label in zip(col_keys, col_labels):
        val = row.get(k, "")
        if isinstance(val, list):
            val = ", ".join(str(v) for v in val)
        val = str(val or "\\u2014")
        lines.append(f"  {{label}}: {{val}}")

lines.append("")
lines.append("-" * 80)
lines.append(f"Total: {{len(DATA)}} items")
lines.append("Generated by LAWA Scouts")

with open("/tmp/report.txt", "w", encoding="utf-8") as f:
    f.write("\\n".join(lines))
print("TXT_GENERATED:/tmp/report.txt")
'''

_PPTX_CODE = '''
import json, textwrap

DATA = {data_json}
TITLE = {title_json}
SUMMARY = {summary_json}
COLUMNS = {columns_json}
ANALYSIS = {analysis_json}

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    import subprocess
    subprocess.check_call(["pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

TEAL = RGBColor(0x0E, 0x74, 0x90)
GRAY = RGBColor(0x47, 0x55, 0x69)
LIGHT_TEAL = RGBColor(0xF0, 0xFD, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x33, 0x41, 0x55)

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = TITLE
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = TEAL
p.alignment = PP_ALIGN.LEFT

p2 = tf.add_paragraph()
p2.text = SUMMARY
p2.font.size = Pt(16)
p2.font.color.rgb = GRAY
p2.space_before = Pt(16)

footer = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.5))
fp = footer.text_frame.paragraphs[0]
fp.text = "LAWA Scouts Report"
fp.font.size = Pt(11)
fp.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)

# Analysis slides (split into pages of ~3 paragraphs each)
if ANALYSIS and len(ANALYSIS.strip()) > 20:
    paragraphs = [p.strip().replace("**", "") for p in ANALYSIS.strip().split("\\n\\n") if p.strip()]
    paras_per_slide = 3
    for slide_start in range(0, len(paragraphs), paras_per_slide):
        slide_paras = paragraphs[slide_start:slide_start + paras_per_slide]
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Heading
        heading_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.6))
        hf = heading_box.text_frame
        hp = hf.paragraphs[0]
        hp.text = "Analysis & Key Findings" if slide_start == 0 else "Analysis (continued)"
        hp.font.size = Pt(20)
        hp.font.bold = True
        hp.font.color.rgb = TEAL

        # Analysis text
        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.5))
        bf = body_box.text_frame
        bf.word_wrap = True
        for pi, para_text in enumerate(slide_paras):
            if pi == 0:
                bp = bf.paragraphs[0]
            else:
                bp = bf.add_paragraph()
            bp.text = textwrap.fill(para_text, width=120)
            bp.font.size = Pt(14)
            bp.font.color.rgb = DARK
            bp.space_after = Pt(14)
            bp.line_spacing = Pt(22)

# Data slide(s) — show up to 15 rows per slide
col_keys = [c["key"] for c in COLUMNS if c.get("type") != "link"]
col_labels = [c["label"] for c in COLUMNS if c.get("type") != "link"]
ncols = len(col_keys)

chunk_size = 15
for chunk_start in range(0, len(DATA), chunk_size):
    chunk = DATA[chunk_start:chunk_start + chunk_size]
    nrows = len(chunk) + 1  # +1 for header

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(12.333)
    height = Inches(6.5)

    table = slide.shapes.add_table(nrows, ncols, left, top, width, height).table

    # Header
    for ci, label in enumerate(col_labels):
        cell = table.cell(0, ci)
        cell.text = label
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.bold = True
            paragraph.font.color.rgb = TEAL
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_TEAL

    # Data rows
    for ri, row in enumerate(chunk, 1):
        for ci, k in enumerate(col_keys):
            val = row.get(k, "")
            if isinstance(val, list):
                val = ", ".join(str(v) for v in val)
            cell = table.cell(ri, ci)
            cell.text = str(val or "\\u2014")[:80]
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.color.rgb = GRAY

prs.save("/tmp/report.pptx")
print("PPTX_GENERATED:/tmp/report.pptx")
'''

# Format registry
_EXPORT_FORMATS = {
    "pdf":   {"template": _PDF_CODE,  "ext": "pdf",  "mime": "application/pdf",
              "pip_deps": ["reportlab"], "uses_summary": True, "uses_topic": True, "uses_analysis": True},
    "excel": {"template": _EXCEL_CODE, "ext": "xlsx", "mime": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
              "pip_deps": [], "uses_summary": False, "uses_topic": False, "uses_analysis": False},
    "html":  {"template": _HTML_CODE,  "ext": "html", "mime": "text/html",
              "pip_deps": [], "uses_summary": True, "uses_topic": False, "uses_analysis": True},
    "csv":   {"template": _CSV_CODE,   "ext": "csv",  "mime": "text/csv",
              "pip_deps": [], "uses_summary": False, "uses_topic": False, "uses_analysis": False},
    "txt":   {"template": _TXT_CODE,   "ext": "txt",  "mime": "text/plain",
              "pip_deps": [], "uses_summary": True, "uses_topic": False, "uses_analysis": True},
    "pptx":  {"template": _PPTX_CODE,  "ext": "pptx", "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
              "pip_deps": ["python-pptx"], "uses_summary": True, "uses_topic": False, "uses_analysis": True},
}


def generate_export(
    data: list[dict],
    columns: list[dict],
    title: str,
    summary: str,
    topic: str,
    export_type: str,
    slug: str | None = None,
    analysis: str = "",
) -> dict:
    """Generate a file in the sandbox for any supported format.

    Returns:
        dict with keys: data (base64), filename, mime_type, error
    """
    if not is_available():
        return {"data": None, "filename": None, "mime_type": None, "error": "E2B not available"}

    fmt = _EXPORT_FORMATS.get(export_type)
    if not fmt:
        return {"data": None, "filename": None, "mime_type": None, "error": f"Unknown format: {export_type}"}

    # Build template kwargs
    kwargs = {
        "data_json": json.dumps(data),
        "columns_json": json.dumps(columns),
        "title_json": json.dumps(title),
    }
    if fmt["uses_summary"]:
        kwargs["summary_json"] = json.dumps(summary)
    if fmt["uses_topic"]:
        kwargs["topic_json"] = json.dumps(topic)
    if fmt.get("uses_analysis"):
        kwargs["analysis_json"] = json.dumps(analysis or "")

    code = fmt["template"].format(**kwargs)
    ext = fmt["ext"]
    mime = fmt["mime"]
    filepath = f"/tmp/report.{ext}"

    try:
        sandbox = Sandbox.create(api_key=get_settings().e2b_api_key)
        try:
            # Install any required packages
            for dep in fmt["pip_deps"]:
                sandbox.run_code(f"import subprocess; subprocess.check_call(['pip', 'install', '{dep}'])")

            execution = sandbox.run_code(code)

            if hasattr(execution, "error") and execution.error:
                logger.warning(f"Export execution error ({export_type}): {execution.error}")
                return {"data": None, "filename": None, "mime_type": None, "error": str(execution.error)}

            # Read the generated file via in-sandbox base64 encoding
            # to avoid binary corruption from sandbox.files.read()
            b64_exec = sandbox.run_code(
                f"import base64\n"
                f"with open('{filepath}', 'rb') as _f:\n"
                f"    print(base64.b64encode(_f.read()).decode(), end='')"
            )
            b64 = ""
            if hasattr(b64_exec, "text") and b64_exec.text:
                b64 = b64_exec.text.strip()
            elif hasattr(b64_exec, "logs") and b64_exec.logs:
                parts = []
                for line in (b64_exec.logs.stdout or []):
                    parts.append(line.text if hasattr(line, "text") else str(line))
                b64 = "".join(parts).strip()

            if b64:
                fname = f"{slug or _slugify(topic or title)}.{ext}"
                logger.info(f"Export generated: {fname} ({len(b64)} chars base64)")
                return {"data": b64, "filename": fname, "mime_type": mime, "error": None}
            else:
                return {"data": None, "filename": None, "mime_type": None, "error": "File not generated"}

        finally:
            try:
                sandbox.kill()
            except Exception:
                pass

    except Exception as e:
        logger.error(f"Export generation failed ({export_type}): {e}")
        return {"data": None, "filename": None, "mime_type": None, "error": str(e)}


# ────────────────────────────────────────────────────────
# Pure-Python Local Export Generation (no E2B needed)
# ────────────────────────────────────────────────────────

# Formats that can be generated locally without E2B sandbox
LOCAL_EXPORTABLE = {"html", "csv", "txt"}


def generate_html_local(
    title: str,
    summary: str,
    columns: list[dict] | None,
    rows: list[dict] | None,
    full_text: str | None,
    slug: str | None = None,
    analysis: str = "",
) -> dict:
    """Generate an HTML export locally, no sandbox needed.
    Works with structured data (columns/rows) OR text-only content.
    """
    import html as html_mod
    import re

    if columns and rows:
        headers_html = "".join(
            f"<th>{html_mod.escape(c.get('label', c.get('key', '')))}</th>"
            for c in columns
        )
        rows_html_parts = []
        for row in rows:
            cells = []
            for c in columns:
                k = c["key"]
                val = row.get(k, "")
                if isinstance(val, list):
                    val = ", ".join(str(v) for v in val)
                val = str(val) if val else "\u2014"
                if c.get("type") == "link" and str(val).startswith("http"):
                    label = html_mod.escape(str(row.get("link_label", "Visit")))
                    cells.append(
                        f'<td><a href="{html_mod.escape(val)}" target="_blank" '
                        f'style="color:#0E7490;text-decoration:underline;">{label}</a></td>'
                    )
                elif c.get("type") == "badge":
                    cells.append(
                        f'<td><span style="display:inline-block;padding:2px 10px;'
                        f'border-radius:999px;background:#ECFEFF;color:#0E7490;'
                        f'font-size:12px;font-weight:600;">{html_mod.escape(val)}</span></td>'
                    )
                else:
                    cells.append(f"<td>{html_mod.escape(val)}</td>")
            rows_html_parts.append(f"<tr>{''.join(cells)}</tr>")

        body_content = f"""
        <div class="summary">{html_mod.escape(summary)}</div>
        <table>
            <thead><tr>{headers_html}</tr></thead>
            <tbody>{''.join(rows_html_parts)}</tbody>
        </table>"""
    else:
        text = full_text or summary or "No content available."
        paragraphs = text.split("\n\n")
        formatted = []
        for p in paragraphs:
            p = p.strip()
            if not p:
                continue
            if p.startswith("---"):
                formatted.append("<hr>")
            else:
                escaped = html_mod.escape(p).replace("\n", "<br>")
                formatted.append(f"<p>{escaped}</p>")
        body_content = f"""
        <div class="summary">{html_mod.escape(summary)}</div>
        <div class="text-content">{''.join(formatted)}</div>"""

    # Build analysis HTML section
    analysis_html = ""
    if analysis and len(analysis.strip()) > 20:
        analysis_paras = []
        for p in analysis.strip().split("\n\n"):
            p = p.strip()
            if p:
                escaped = html_mod.escape(p)
                # Convert **bold** to <strong>
                escaped = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', escaped)
                escaped = escaped.replace("\n", "<br>")
                analysis_paras.append(f"<p>{escaped}</p>")
        if analysis_paras:
            analysis_html = f"""
        <div class="analysis">
            <h2>Analysis &amp; Key Findings</h2>
            {''.join(analysis_paras)}
        </div>"""

    page_html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>{html_mod.escape(title)}</title>
<style>
  *{{margin:0;padding:0;box-sizing:border-box}}
  body{{font-family:'Inter','Segoe UI',system-ui,sans-serif;background:#F8FAFC;color:#334155;padding:40px 24px}}
  .container{{max-width:1100px;margin:0 auto}}
  h1{{font-size:24px;font-weight:700;color:#0E7490;margin-bottom:8px}}
  .subtitle{{font-size:13px;color:#64748B;margin-bottom:24px}}
  .summary{{font-size:14px;line-height:1.7;color:#475569;margin-bottom:24px;padding:16px 20px;background:#fff;border-radius:12px;border:1px solid #E2E8F0}}
  .analysis{{margin-bottom:28px;padding:20px 24px;background:#fff;border-radius:12px;border:1px solid #E2E8F0}}
  .analysis h2{{font-size:16px;font-weight:700;color:#0E7490;margin-bottom:14px}}
  .analysis p{{font-size:14px;line-height:1.75;color:#475569;margin-bottom:12px}}
  .analysis p:last-child{{margin-bottom:0}}
  .text-content{{font-size:14px;line-height:1.8;color:#475569}}
  .text-content p{{margin-bottom:16px}}
  .text-content hr{{border:none;border-top:1px solid #E2E8F0;margin:24px 0}}
  h3{{font-size:14px;font-weight:600;color:#0E7490;margin:24px 0 8px}}
  table{{width:100%;border-collapse:collapse;background:#fff;border-radius:12px;overflow:hidden;border:1px solid #E2E8F0}}
  th{{background:#F0FDFA;color:#0E7490;font-size:12px;font-weight:600;text-transform:uppercase;letter-spacing:.5px;padding:12px 16px;text-align:left;border-bottom:2px solid #E2E8F0}}
  td{{padding:10px 16px;font-size:13px;border-bottom:1px solid #F1F5F9;vertical-align:top}}
  tr:hover td{{background:#F8FAFC}}
  .footer{{margin-top:32px;text-align:center;font-size:11px;color:#94A3B8}}
</style>
</head>
<body>
<div class="container">
  <h1>{html_mod.escape(title)}</h1>
  <p class="subtitle">LAWA Scouts Report</p>
  {body_content}
  {analysis_html}
  <p class="footer">Generated by LAWA Scouts</p>
</div>
</body>
</html>"""

    b64 = base64.b64encode(page_html.encode("utf-8")).decode()
    fname = f"{slug or _slugify(title)}.html"
    return {"data": b64, "filename": fname, "mime_type": "text/html", "error": None}


def generate_csv_local(
    columns: list[dict] | None,
    rows: list[dict] | None,
    full_text: str | None,
    slug: str | None = None,
    title: str = "Report",
) -> dict:
    """Generate a CSV export locally, no sandbox needed."""
    import csv
    import io

    output = io.StringIO()
    writer = csv.writer(output)

    if columns and rows:
        col_keys = [c["key"] for c in columns]
        col_labels = [c.get("label", c["key"]) for c in columns]
        writer.writerow(col_labels)
        for row in rows:
            cells = []
            for k in col_keys:
                val = row.get(k, "")
                if isinstance(val, list):
                    val = ", ".join(str(v) for v in val)
                cells.append(str(val or ""))
            writer.writerow(cells)
    else:
        writer.writerow(["Content"])
        for line in (full_text or "").split("\n"):
            writer.writerow([line])

    csv_bytes = output.getvalue().encode("utf-8")
    b64 = base64.b64encode(csv_bytes).decode()
    fname = f"{slug or _slugify(title)}.csv"
    return {"data": b64, "filename": fname, "mime_type": "text/csv", "error": None}


def generate_txt_local(
    title: str,
    summary: str,
    columns: list[dict] | None,
    rows: list[dict] | None,
    full_text: str | None,
    slug: str | None = None,
    analysis: str = "",
) -> dict:
    """Generate a plain text export locally, no sandbox needed."""
    import textwrap

    lines = []
    lines.append("=" * 80)
    lines.append(title)
    lines.append("=" * 80)
    lines.append("")
    lines.append(textwrap.fill(summary, width=80))
    lines.append("")

    # Analysis section
    if analysis and len(analysis.strip()) > 20:
        lines.append("-" * 80)
        lines.append("ANALYSIS & KEY FINDINGS")
        lines.append("-" * 80)
        lines.append("")
        for para in analysis.strip().split("\n\n"):
            para = para.strip().replace("**", "")
            if para:
                lines.append(textwrap.fill(para, width=80))
                lines.append("")

    if columns and rows:
        lines.append("-" * 80)
        lines.append("DETAILED DATA")
        lines.append("-" * 80)
        col_keys = [c["key"] for c in columns if c.get("type") != "link"]
        col_labels = [c.get("label", c["key"]) for c in columns if c.get("type") != "link"]
        for ri, row in enumerate(rows, 1):
            lines.append(f"\n[{ri}]")
            for k, label in zip(col_keys, col_labels):
                val = row.get(k, "")
                if isinstance(val, list):
                    val = ", ".join(str(v) for v in val)
                val = str(val) if val else "\u2014"
                lines.append(f"  {label}: {val}")
        lines.append("")
        lines.append("-" * 80)
        lines.append(f"Total: {len(rows)} items")
    else:
        lines.append("-" * 80)
        lines.append(full_text or "No content available.")

    lines.append("")
    lines.append("Generated by LAWA Scouts")

    txt_content = "\n".join(lines)
    b64 = base64.b64encode(txt_content.encode("utf-8")).decode()
    fname = f"{slug or _slugify(title)}.txt"
    return {"data": b64, "filename": fname, "mime_type": "text/plain", "error": None}


def generate_local_export(
    export_type: str,
    title: str,
    summary: str,
    columns: list[dict] | None = None,
    rows: list[dict] | None = None,
    full_text: str | None = None,
    slug: str | None = None,
    analysis: str = "",
) -> dict:
    """Generate an export using pure Python (no E2B sandbox).
    Supports: html, csv, txt. Returns same dict shape as generate_export().
    """
    try:
        if export_type == "html":
            return generate_html_local(title, summary, columns, rows, full_text, slug, analysis)
        elif export_type == "csv":
            return generate_csv_local(columns, rows, full_text, slug, title)
        elif export_type == "txt":
            return generate_txt_local(title, summary, columns, rows, full_text, slug, analysis)
        else:
            return {"data": None, "filename": None, "mime_type": None,
                    "error": f"Local generation not supported for: {export_type}"}
    except Exception as e:
        logger.error(f"Local export generation failed ({export_type}): {e}")
        return {"data": None, "filename": None, "mime_type": None, "error": str(e)}


# ────────────────────────────────────────────────────────
# Shared sandbox execution
# ────────────────────────────────────────────────────────

def _execute_in_sandbox(full_code: str) -> dict:
    """Run code in E2B sandbox and collect stdout + chart PNGs."""
    sandbox = Sandbox.create(api_key=get_settings().e2b_api_key)
    try:
        execution = sandbox.run_code(full_code)

        analysis_text = ""
        if hasattr(execution, "text"):
            analysis_text = execution.text or ""
        elif hasattr(execution, "logs") and execution.logs:
            analysis_text = "\n".join(
                line.text for line in (execution.logs.stdout or [])
                if hasattr(line, "text")
            )

        # If code decided data isn't chartable
        if "NO_CHARTS" in analysis_text:
            logger.info("Sandbox: code decided data is not chartable")
            return {"analysis_text": "", "charts": [], "error": None}

        charts = []
        if hasattr(execution, "results"):
            for result in execution.results:
                if hasattr(result, "png") and result.png:
                    charts.append(result.png)

        if not charts:
            for i in range(1, 4):
                try:
                    content = sandbox.files.read(f"/tmp/chart_{i}.png")
                    if content:
                        charts.append(base64.b64encode(
                            content if isinstance(content, bytes) else content.encode()
                        ).decode())
                except Exception:
                    break

        error = None
        if hasattr(execution, "error") and execution.error:
            error = str(execution.error)
            logger.warning(f"Sandbox execution error: {error}")

        logger.info(f"Sandbox: analysis={len(analysis_text)} chars, charts={len(charts)}")
        return {"analysis_text": analysis_text, "charts": charts, "error": error}

    finally:
        try:
            sandbox.kill()
        except Exception:
            pass
